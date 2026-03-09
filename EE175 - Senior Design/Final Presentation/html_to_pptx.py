import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation

def convert_to_pptx(html_path, out_pptx):
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'html.parser')

    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    slides = soup.find_all('section', class_='slide')
    
    for idx, slide_tag in enumerate(slides):
        if idx == 0:
            slide = prs.slides.add_slide(title_slide_layout)
            title_text = slide_tag.find('h1')
            if title_text:
                slide.shapes.title.text = title_text.get_text(separator=" ", strip=True)
                
            lead_text = slide_tag.find(class_='lead')
            if lead_text:
                slide.placeholders[1].text = lead_text.get_text(separator=" ", strip=True)
            continue
            
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        h2 = slide_tag.find('h2')
        if h2:
            slide.shapes.title.text = h2.get_text(separator=" ", strip=True)
            
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        content_div = slide_tag.find(class_='content')
        if not content_div:
            continue
            
        subhead = content_div.find(class_='subhead')
        if subhead:
            p = text_frame.paragraphs[0]
            p.text = subhead.get_text(separator=" ", strip=True)
            p.level = 0

        # Gather relevant containers
        elements_to_parse = []
        for elem in content_div.descendants:
            if getattr(elem, 'name', None) in ['h3', 'h4', 'strong', 'li', 'p', 'th', 'td', 'span']:
                # skip layout elements
                if 'subhead' in elem.get('class', []): continue
                # Do not parse spans if they are inside a component we already parse, but let's just parse all text.
                if elem.name == 'span' and 'status-pill' not in elem.get('class', []): continue
                
                # Simple dedup for nested text (like strong inside p)
                if elem.name in ['h3', 'h4', 'strong', 'th']:
                    level = 0
                else:
                    level = 1
                
                text = elem.get_text(separator=" ", strip=True)
                if text and elem.find_parents(['h3', 'h4', 'strong', 'th']) == []: # ignore child nodes of already parsed headings
                     elements_to_parse.append((text, level))

        # Filter out exact duplicates sequentially
        filtered_elements = []
        for text, level in elements_to_parse:
             # deduplicate string from the subhead
             if subhead and text == subhead.get_text(separator=" ", strip=True):
                 continue
             
             if not filtered_elements or filtered_elements[-1][0] != text:
                 # Check if the text is entirely contained in the parent
                 if filtered_elements:
                     prev_text, prev_lvl = filtered_elements[-1]
                     if text in prev_text and len(text) > 0 and prev_lvl == level:
                         continue
                 filtered_elements.append((text, level))

        for (text, level) in filtered_elements:
             if len(text_frame.paragraphs) == 1 and text_frame.paragraphs[0].text == "":
                 p = text_frame.paragraphs[0]
             else:
                 p = text_frame.add_paragraph()
             p.text = text
             p.level = level
            
    prs.save(out_pptx)
    print("Created", out_pptx)

if __name__ == '__main__':
    convert_to_pptx('slides.html', 'BMS_Final_Design_Presentation_Editable.pptx')

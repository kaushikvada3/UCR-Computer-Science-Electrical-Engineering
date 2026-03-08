import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches
import os
import sys

async def generate_pptx(html_path, output_pptx):
    # Ensure correct path format for Windows
    file_url = f"file:///{html_path.replace(chr(92), '/')}"
    
    print(f"Loading {file_url}...")
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        # Large viewport to ensure all slides are rendered properly.
        # It's a grid, so let's make it tall to prevent any scrolling issues.
        page = await browser.new_page(viewport={'width': 1400, 'height': 8500})
        await page.goto(file_url, wait_until="networkidle")
        
        # Wait a second to allow fonts and gradients to settle
        await page.wait_for_timeout(2000)
        
        prs = Presentation()
        # Set slide size to match 1280x720 aspect ratio (16:9)
        prs.slide_width = Inches(13.3333333333)
        prs.slide_height = Inches(7.5)
        
        # Blank slide layout is index 6
        blank_slide_layout = prs.slide_layouts[6]
        
        for i in range(1, 11):
            slide_elem = await page.query_selector(f"#slide{i}")
            if slide_elem:
                print(f"Capturing slide {i}...")
                screenshot_path = f"slide_{i}.png"
                # The slide elements are 1280x720
                await slide_elem.screenshot(path=screenshot_path)
                
                # Add to PPTX
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(screenshot_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # Cleanup screenshot
                os.remove(screenshot_path)
            else:
                print(f"Slide {i} not found.")
        
        await browser.close()
        
        prs.save(output_pptx)
        print(f"\nSuccessfully saved presentation to {output_pptx}")

if __name__ == "__main__":
    if sys.platform == "win32":
        asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())
    script_dir = os.path.dirname(os.path.abspath(__file__))
    html_file = os.path.join(script_dir, "slides.html")
    pptx_file = os.path.join(script_dir, "BMS_Final_Design_Presentation.pptx")
    asyncio.run(generate_pptx(html_file, pptx_file))
